import os
import glob
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

def extract_text_from_pptx(filepath):
    text = []
    if not Presentation:
        return "python-pptx not installed"
    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
    except Exception as e:
        return f"Error reading {filepath}: {e}"
    return "\n".join(text)

def main():
    folder = "."
    files = glob.glob(os.path.join(folder, "*.pptx"))
    with open("extracted_text.txt", "w", encoding="utf-8") as out:
        for file in files:
            out.write(f"--- {os.path.basename(file)} ---\n")
            out.write(extract_text_from_pptx(file)[:1000] + "...\n\n")

if __name__ == "__main__":
    main()
